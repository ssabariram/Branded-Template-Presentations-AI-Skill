"""
<company_name> Template Injector — Generates presentations from a structured content plan.

Usage:
    python <company_name>_injector.py <content_plan.json> [--template <path>] [--output <path>]

The content plan is a JSON file with the following structure:
{
  "slides": [
    {
      "type": "COVER",
      "title": "Presentation Title",
      "subtitle": "Optional subtitle"
    },
    {
      "type": "CONTENT",
      "title": "Slide Title",
      "bullets": [
        "First bullet point",
        "Second bullet point",
        {"text": "Bullet with sub-bullets", "sub_bullets": ["Sub 1", "Sub 2"]}
      ]
    },
    ...
  ]
}
"""

import argparse
import io
import json
import math
import sys
import copy
from pathlib import Path

# Slide transition support (raw lxml — python-pptx has no public API for this)
try:
    from transitions import inject_transition, validate_transition, TRANSITION_CATALOG, RECOMMENDED_TRANSITION
    _TRANSITIONS_AVAILABLE = True
except ImportError:
    _TRANSITIONS_AVAILABLE = False

# Element animation support (raw lxml — python-pptx has no public API for this)
try:
    from animations import inject_animations, validate_animations, ANIMATION_CATALOG, RECOMMENDED_ANIMATIONS
    _ANIMATIONS_AVAILABLE = True
except ImportError:
    _ANIMATIONS_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
except ImportError:
    print("ERROR: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)


# ─── Brand Colors & Chart Constants ───

<company_name>_CHART_COLORS = [
    RGBColor(0x00, 0x80, 0xBA),  # Blue
    RGBColor(0xFF, 0x99, 0x01),  # Orange
    RGBColor(0x00, 0x44, 0x7A),  # Dark Blue
    RGBColor(0x4C, 0xAF, 0x50),  # Green
    RGBColor(0x5B, 0xC0, 0xDE),  # Light Blue
    RGBColor(0x98, 0x99, 0x98),  # Grey
]

CHART_TYPE_MAP = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
    "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
    "area": XL_CHART_TYPE.AREA,
}

# Body area coordinates from Layout 14 (Title and Content)
BODY_AREA = {
    "left": Inches(0.76),
    "top": Inches(1.52),
    "width": Inches(12.09),
    "height": Inches(4.76),
}

# Split layout: left half for bullets, right half for chart/table
SPLIT_LEFT = {
    "left": Inches(0.76),
    "top": Inches(1.52),
    "width": Inches(5.30),
    "height": Inches(4.76),
}

SPLIT_RIGHT = {
    "left": Inches(6.26),
    "top": Inches(1.52),
    "width": Inches(5.90),
    "height": Inches(4.76),
}


def _parse_hex_color(hex_str):
    """Parse a hex color string like '#0080BA' or '0080BA' into RGBColor."""
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _remove_body_placeholder(slide):
    """Remove the body placeholder (idx=1) from a slide's shape tree."""
    from pptx.oxml.ns import qn
    spTree = slide._element.find(qn('p:cSld')).find(qn('p:spTree'))
    for sp in list(spTree):
        tag = sp.tag.split('}')[-1] if '}' in sp.tag else sp.tag
        if tag == 'sp':
            nvSpPr = sp.find(qn('p:nvSpPr'))
            if nvSpPr is not None:
                nvPr = nvSpPr.find(qn('p:nvPr'))
                if nvPr is not None:
                    ph = nvPr.find(qn('p:ph'))
                    if ph is not None and ph.get('idx') == '1':
                        spTree.remove(sp)
                        return True
    return False


def _remove_placeholder_by_idx(slide, idx):
    """Remove a placeholder by its idx attribute from the slide's shape tree."""
    from pptx.oxml.ns import qn
    spTree = slide._element.find(qn('p:cSld')).find(qn('p:spTree'))
    for sp in list(spTree):
        tag = sp.tag.split('}')[-1] if '}' in sp.tag else sp.tag
        if tag == 'sp':
            nvSpPr = sp.find(qn('p:nvSpPr'))
            if nvSpPr is not None:
                nvPr = nvSpPr.find(qn('p:nvPr'))
                if nvPr is not None:
                    ph = nvPr.find(qn('p:ph'))
                    if ph is not None and ph.get('idx') == str(idx):
                        spTree.remove(sp)
                        return True
    return False


# ─── Layout Catalog (hardcoded from manual analysis) ───

LAYOUT_CATALOG = {
    "COVER": {
        "layout_index": 0,
        "layout_name": "8_Title Slide 3",
        "description": "Dark blue bg + city skyline. Opening slide.",
        "placeholders": {
            "title": {"idx": 0, "max_chars_per_line": 20, "max_lines": 2},
            "subtitle": {"idx": 1, "max_chars_per_line": 35, "max_lines": 2},
        }
    },
    "COVER_ALT": {
        "layout_index": 8,
        "layout_name": "2_Title Slide",
        "description": "Blue bg + people photo. Alt opening slide.",
        "placeholders": {
            "title": {"idx": 0, "max_chars_per_line": 20, "max_lines": 3},
            "subtitle": {"idx": 1, "max_chars_per_line": 45, "max_lines": 2},
        }
    },
    "COVER_FULL": {
        "layout_index": 6,
        "layout_name": "Title Slide 2 Blue",
        "description": "Full blue bg, bold statement cover.",
        "placeholders": {
            "title": {"idx": 0, "max_chars_per_line": 30, "max_lines": 2},
            "subtitle": {"idx": 1, "max_chars_per_line": 45, "max_lines": 2},
        }
    },
    "CHAPTER": {
        "layout_index": 24,
        "layout_name": "2_Chapter Slide",
        "description": "White bg + city photo corner. Chapter divider.",
        "placeholders": {
            "title": {"idx": 0, "max_chars_per_line": 18, "max_lines": 2},
            "subtitle": {"idx": 1, "max_chars_per_line": 35, "max_lines": 1},
        }
    },
    "SECTION_BLUE": {
        "layout_index": 23,
        "layout_name": "2_Section Slide",
        "description": "Blue band section divider.",
        "placeholders": {
            "title": {"idx": 0, "max_chars_per_line": 17, "max_lines": 2},
        }
    },
    "SECTION_GREY": {
        "layout_index": 22,
        "layout_name": "1_Section Slide",
        "description": "Grey band section divider.",
        "placeholders": {
            "title": {"idx": 0, "max_chars_per_line": 17, "max_lines": 2},
        }
    },
    "CONTENT": {
        "layout_index": 14,
        "layout_name": "Title and Content",
        "description": "Primary content slide with bullets.",
        "placeholders": {
            "title": {"idx": 0, "max_chars_per_line": 30, "max_lines": 1},
            "body": {"idx": 1, "max_chars_per_line": 55, "max_lines": 7, "max_bullets": 6},
        }
    },
    "CONTENT_SIDEBAR": {
        "layout_index": 15,
        "layout_name": "1_Title and Content",
        "description": "Content with sidebar accent. More space.",
        "placeholders": {
            "title": {"idx": 0, "max_chars_per_line": 15, "max_lines": 1},
            "body": {"idx": 1, "max_chars_per_line": 55, "max_lines": 10, "max_bullets": 8},
        }
    },
    "TWO_COLUMN": {
        "layout_index": 16,
        "layout_name": "Two Content",
        "description": "Two columns side by side.",
        "placeholders": {
            "title": {"idx": 0, "max_chars_per_line": 30, "max_lines": 1},
            "left": {"idx": 1, "max_chars_per_line": 25, "max_lines": 7, "max_bullets": 5},
            "right": {"idx": 2, "max_chars_per_line": 28, "max_lines": 7, "max_bullets": 5},
        }
    },
    "QUOTE": {
        "layout_index": 25,
        "layout_name": "1_Large Quote",
        "description": "Large centered text for quotes/stats.",
        "placeholders": {
            "title": {"idx": 0, "max_chars_per_line": 30, "max_lines": 3},
        }
    },
    "CLOSING": {
        "layout_index": 26,
        "layout_name": "2_Title Slide 3 City End",
        "description": "Thank You / closing slide. No editable placeholders.",
        "placeholders": {}
    },
    # ─── New slide types (all use Layout 14: Title and Content) ───
    "CHART": {
        "layout_index": 14,
        "layout_name": "Title and Content",
        "description": "Full-width chart replacing body area.",
        "placeholders": {
            "title": {"idx": 0, "max_chars_per_line": 30, "max_lines": 1},
        }
    },
    "TABLE": {
        "layout_index": 14,
        "layout_name": "Title and Content",
        "description": "Full-width styled table replacing body area.",
        "placeholders": {
            "title": {"idx": 0, "max_chars_per_line": 30, "max_lines": 1},
        }
    },
    "CONTENT_CHART": {
        "layout_index": 14,
        "layout_name": "Title and Content",
        "description": "Split: bullets left, chart right.",
        "placeholders": {
            "title": {"idx": 0, "max_chars_per_line": 30, "max_lines": 1},
            "body": {"idx": 1, "max_chars_per_line": 40, "max_lines": 7, "max_bullets": 6},
        }
    },
    "CONTENT_TABLE": {
        "layout_index": 14,
        "layout_name": "Title and Content",
        "description": "Split: bullets left, table right.",
        "placeholders": {
            "title": {"idx": 0, "max_chars_per_line": 30, "max_lines": 1},
            "body": {"idx": 1, "max_chars_per_line": 40, "max_lines": 7, "max_bullets": 6},
        }
    },
    "KPI": {
        "layout_index": 14,
        "layout_name": "Title and Content",
        "description": "2-4 metric cards as large rounded rectangles.",
        "placeholders": {
            "title": {"idx": 0, "max_chars_per_line": 30, "max_lines": 1},
        }
    },
    "TIMELINE": {
        "layout_index": 14,
        "layout_name": "Title and Content",
        "description": "Horizontal timeline with 2-6 milestone nodes.",
        "placeholders": {
            "title": {"idx": 0, "max_chars_per_line": 30, "max_lines": 1},
        }
    },
}


# ─── Validation ───

def validate_text(text, max_chars_per_line, max_lines, field_name, slide_num):
    """Validate text fits within constraints. Returns warnings list."""
    warnings = []
    if text is None:
        return warnings

    lines = text.split('\n')
    if len(lines) > max_lines:
        warnings.append(
            f"Slide {slide_num}, {field_name}: {len(lines)} lines exceeds max {max_lines}"
        )

    for i, line in enumerate(lines):
        if len(line) > max_chars_per_line:
            warnings.append(
                f"Slide {slide_num}, {field_name}, line {i+1}: {len(line)} chars exceeds max {max_chars_per_line} "
                f"(\"{line[:30]}...\")"
            )

    return warnings


def validate_bullets(bullets, max_chars_per_line, max_bullets, field_name, slide_num):
    """Validate bullet list fits within constraints."""
    warnings = []
    if bullets is None:
        return warnings

    if len(bullets) > max_bullets:
        warnings.append(
            f"Slide {slide_num}, {field_name}: {len(bullets)} bullets exceeds max {max_bullets}"
        )

    for i, bullet in enumerate(bullets):
        text = bullet if isinstance(bullet, str) else bullet.get("text", "")
        if len(text) > max_chars_per_line:
            warnings.append(
                f"Slide {slide_num}, {field_name}, bullet {i+1}: {len(text)} chars exceeds max {max_chars_per_line} "
                f"(\"{text[:30]}...\")"
            )

        # Check sub-bullets
        if isinstance(bullet, dict) and "sub_bullets" in bullet:
            for j, sub in enumerate(bullet["sub_bullets"]):
                if len(sub) > max_chars_per_line:
                    warnings.append(
                        f"Slide {slide_num}, {field_name}, bullet {i+1} sub {j+1}: "
                        f"{len(sub)} chars exceeds max {max_chars_per_line}"
                    )

    return warnings


def validate_chart(chart_spec, field_name, slide_num):
    """Validate a chart specification."""
    warnings = []
    if not chart_spec:
        return warnings

    chart_type = chart_spec.get("type", "").lower()
    valid_types = list(CHART_TYPE_MAP.keys())
    if chart_type not in valid_types:
        warnings.append(
            f"Slide {slide_num}, {field_name}: Unknown chart type '{chart_type}'. "
            f"Valid types: {', '.join(valid_types)}"
        )

    categories = chart_spec.get("categories", [])
    series_list = chart_spec.get("series", [])

    if not categories:
        warnings.append(f"Slide {slide_num}, {field_name}: Chart has no categories")
    if not series_list:
        warnings.append(f"Slide {slide_num}, {field_name}: Chart has no series")

    for s_idx, series in enumerate(series_list):
        values = series.get("values", [])
        if len(values) != len(categories):
            warnings.append(
                f"Slide {slide_num}, {field_name}, series {s_idx+1}: "
                f"{len(values)} values but {len(categories)} categories"
            )

    if chart_type in ("pie", "doughnut") and len(categories) > 8:
        warnings.append(
            f"Slide {slide_num}, {field_name}: Pie/doughnut charts recommend max 8 categories "
            f"(has {len(categories)})"
        )

    return warnings


def validate_table(table_spec, field_name, slide_num):
    """Validate a table specification."""
    warnings = []
    if not table_spec:
        return warnings

    headers = table_spec.get("headers", [])
    rows = table_spec.get("rows", [])
    n_cols = len(headers)

    if not headers:
        warnings.append(f"Slide {slide_num}, {field_name}: Table has no headers")

    if n_cols > 6:
        warnings.append(
            f"Slide {slide_num}, {field_name}: {n_cols} columns exceeds recommended max of 6"
        )

    if len(rows) > 10:
        warnings.append(
            f"Slide {slide_num}, {field_name}: {len(rows)} rows exceeds recommended max of 10"
        )

    for r_idx, row in enumerate(rows):
        if len(row) != n_cols:
            warnings.append(
                f"Slide {slide_num}, {field_name}, row {r_idx+1}: "
                f"{len(row)} cells but {n_cols} headers"
            )

    return warnings


def validate_kpis(kpis, slide_num):
    """Validate KPI card specifications."""
    warnings = []
    if not kpis:
        return warnings

    if len(kpis) < 2:
        warnings.append(f"Slide {slide_num}: KPI slide recommends at least 2 cards (has {len(kpis)})")
    if len(kpis) > 4:
        warnings.append(f"Slide {slide_num}: KPI slide recommends max 4 cards (has {len(kpis)})")

    for k_idx, kpi in enumerate(kpis):
        value = kpi.get("value", "")
        label = kpi.get("label", "")
        if len(value) > 8:
            warnings.append(
                f"Slide {slide_num}, KPI {k_idx+1}: value '{value}' exceeds 8 chars"
            )
        if len(label) > 20:
            warnings.append(
                f"Slide {slide_num}, KPI {k_idx+1}: label '{label}' exceeds 20 chars"
            )

    return warnings


def validate_milestones(milestones, slide_num):
    """Validate timeline milestone specifications."""
    warnings = []
    if not milestones:
        return warnings

    if len(milestones) < 2:
        warnings.append(f"Slide {slide_num}: Timeline recommends at least 2 milestones (has {len(milestones)})")
    if len(milestones) > 6:
        warnings.append(f"Slide {slide_num}: Timeline recommends max 6 milestones (has {len(milestones)})")

    for m_idx, ms in enumerate(milestones):
        if not ms.get("label"):
            warnings.append(f"Slide {slide_num}, milestone {m_idx+1}: missing required 'label' field")

    return warnings


def validate_plan(plan):
    """Validate entire content plan against layout contracts."""
    all_warnings = []

    for i, slide in enumerate(plan["slides"], 1):
        slide_type = slide.get("type", "").upper()

        if slide_type not in LAYOUT_CATALOG:
            all_warnings.append(f"Slide {i}: Unknown type '{slide_type}'")
            continue

        catalog = LAYOUT_CATALOG[slide_type]
        phs = catalog["placeholders"]

        # Validate title
        if "title" in phs and "title" in slide:
            ph = phs["title"]
            all_warnings.extend(
                validate_text(slide["title"], ph["max_chars_per_line"], ph["max_lines"], "title", i)
            )

        # Validate subtitle
        if "subtitle" in phs and "subtitle" in slide:
            ph = phs["subtitle"]
            all_warnings.extend(
                validate_text(slide["subtitle"], ph["max_chars_per_line"], ph["max_lines"], "subtitle", i)
            )

        # Validate body bullets
        if "body" in phs and "bullets" in slide:
            ph = phs["body"]
            all_warnings.extend(
                validate_bullets(slide["bullets"], ph["max_chars_per_line"], ph.get("max_bullets", 99), "body", i)
            )

        # Validate two-column
        if "left" in phs and "left_bullets" in slide:
            ph = phs["left"]
            all_warnings.extend(
                validate_bullets(slide["left_bullets"], ph["max_chars_per_line"], ph.get("max_bullets", 99), "left", i)
            )
        if "right" in phs and "right_bullets" in slide:
            ph = phs["right"]
            all_warnings.extend(
                validate_bullets(slide["right_bullets"], ph["max_chars_per_line"], ph.get("max_bullets", 99), "right", i)
            )

        # Validate chart specifications
        if "chart" in slide:
            all_warnings.extend(validate_chart(slide["chart"], "chart", i))
        if "left_chart" in slide:
            all_warnings.extend(validate_chart(slide["left_chart"], "left_chart", i))
        if "right_chart" in slide:
            all_warnings.extend(validate_chart(slide["right_chart"], "right_chart", i))

        # Validate table specifications
        if "table" in slide:
            all_warnings.extend(validate_table(slide["table"], "table", i))

        # Validate KPIs
        if slide_type == "KPI" and "kpis" in slide:
            all_warnings.extend(validate_kpis(slide["kpis"], i))

        # Validate timeline milestones
        if slide_type == "TIMELINE" and "milestones" in slide:
            all_warnings.extend(validate_milestones(slide["milestones"], i))

        # Mutual exclusion warnings
        if slide_type == "CONTENT" and "chart" in slide and "bullets" in slide:
            all_warnings.append(
                f"Slide {i}: CONTENT has both 'chart' and 'bullets'. "
                f"Chart will be ignored; use CONTENT_CHART for side-by-side."
            )
        if slide_type == "TWO_COLUMN":
            if "left_chart" in slide and "left_bullets" in slide:
                all_warnings.append(
                    f"Slide {i}: TWO_COLUMN has both 'left_chart' and 'left_bullets'. "
                    f"left_bullets will render; left_chart will replace the placeholder."
                )
            if "right_chart" in slide and "right_bullets" in slide:
                all_warnings.append(
                    f"Slide {i}: TWO_COLUMN has both 'right_chart' and 'right_bullets'. "
                    f"right_bullets will render; right_chart will replace the placeholder."
                )

        # Validate icons
        if "icons" in slide:
            for ic_idx, ic in enumerate(slide["icons"]):
                icon_name = ic.get("name", "").lower().replace(" ", "_").replace("-", "_")
                if icon_name and icon_name not in ICON_CATALOG:
                    all_warnings.append(
                        f"Slide {i}, icon {ic_idx+1}: Unknown icon '{icon_name}'"
                    )

        # Validate transition
        if "transition" in slide and _TRANSITIONS_AVAILABLE:
            all_warnings.extend(validate_transition(slide["transition"], i))

        # Validate animations
        if "animations" in slide and _ANIMATIONS_AVAILABLE:
            all_warnings.extend(validate_animations(slide["animations"], i))

    return all_warnings


# ─── Injection ───

def clear_placeholder(placeholder):
    """Clear all text from a placeholder while preserving formatting of first paragraph."""
    tf = placeholder.text_frame
    # Keep the first paragraph (it carries the master formatting)
    # but clear its text
    for i, para in enumerate(tf.paragraphs):
        if i == 0:
            # Clear runs in the first paragraph
            for run in para.runs:
                run.text = ""
        # We can't easily remove extra paragraphs with python-pptx,
        # so we just clear them
        else:
            for run in para.runs:
                run.text = ""


def set_placeholder_text(placeholder, text):
    """Set text in a placeholder, preserving the existing formatting."""
    tf = placeholder.text_frame

    # If there are runs, use the first run's formatting
    if tf.paragraphs and tf.paragraphs[0].runs:
        first_run = tf.paragraphs[0].runs[0]
        first_run.text = text
        # Clear any extra runs
        for run in tf.paragraphs[0].runs[1:]:
            run.text = ""
        # Clear extra paragraphs
        for para in tf.paragraphs[1:]:
            for run in para.runs:
                run.text = ""
    else:
        # Fallback: just set the text directly
        tf.paragraphs[0].text = text


def set_bullets(placeholder, bullets):
    """
    Set bullet content in a body placeholder.

    bullets can be:
    - A list of strings: ["bullet 1", "bullet 2"]
    - A list of dicts: [{"text": "bullet", "sub_bullets": ["sub1", "sub2"]}]
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    tf = placeholder.text_frame
    txBody = tf._txBody

    # Grab the first paragraph's XML as a formatting template
    first_para = tf.paragraphs[0]
    para_template = copy.deepcopy(first_para._p)

    # Remove all existing paragraphs
    for p in txBody.findall(qn('a:p')):
        txBody.remove(p)

    for bullet in bullets:
        if isinstance(bullet, str):
            # Simple bullet
            new_p = copy.deepcopy(para_template)
            _set_paragraph_text(new_p, bullet)
            txBody.append(new_p)

        elif isinstance(bullet, dict):
            # Bullet with potential sub-bullets
            new_p = copy.deepcopy(para_template)
            _set_paragraph_text(new_p, bullet.get("text", ""))
            txBody.append(new_p)

            # Add sub-bullets at indent level 1
            for sub_text in bullet.get("sub_bullets", []):
                sub_p = copy.deepcopy(para_template)
                _set_paragraph_text(sub_p, sub_text)
                # Set indent level to 1
                pPr = sub_p.find(qn('a:pPr'))
                if pPr is None:
                    pPr = etree.SubElement(sub_p, qn('a:pPr'))
                pPr.set('lvl', '1')
                txBody.append(sub_p)


def _set_paragraph_text(p_elem, text):
    """Set the text content of a paragraph XML element."""
    from pptx.oxml.ns import qn

    # Find existing runs and set text in the first one
    runs = p_elem.findall(qn('a:r'))
    if runs:
        # Set text in first run
        rText = runs[0].find(qn('a:t'))
        if rText is not None:
            rText.text = text
        else:
            rText = etree.SubElement(runs[0], qn('a:t'))
            rText.text = text

        # Remove extra runs
        for r in runs[1:]:
            p_elem.remove(r)
    else:
        # No runs exist — create one
        from lxml import etree
        r = etree.SubElement(p_elem, qn('a:r'))
        rPr = etree.SubElement(r, qn('a:rPr'))
        rPr.set('lang', 'en-US')
        t = etree.SubElement(r, qn('a:t'))
        t.text = text


# ─── Icon Catalog (extracted from template media library) ───

ICON_CATALOG = {
    # Slide 26 — Solution Category Icons
    "managed_services_wheel": {"png": "image48.png", "svg": "image49.svg"},
    "cyber_security": {"png": "image50.png", "svg": "image51.svg"},
    "digital": {"png": "image52.png", "svg": "image53.svg"},
    "modern_networking": {"png": "image54.png", "svg": "image55.svg"},
    "modern_platforms": {"png": "image56.png", "svg": "image57.svg"},
    "total_experience": {"png": "image58.png", "svg": "image59.svg"},
    "cost_optimization": {"png": "image60.png", "svg": "image61.svg"},
    # Slide 27 — Brand Icon Library page 1
    "cloud": {"png": "image62.png", "svg": "image63.svg"},
    "infrastructure_modernization": {"png": "image65.png", "svg": "image66.svg"},
    "workforce_transformation": {"png": "image67.png", "svg": "image68.svg"},
    "cybersecurity": {"png": "image50.png", "svg": "image64.svg"},
    "lifecycle_services": {"png": "image115.png", "svg": "image116.svg"},
    "digital_transformation": {"png": "image69.png", "svg": "image70.svg"},
    "application_stack": {"png": "image73.png", "svg": "image74.svg"},
    "application_development": {"png": "image71.png", "svg": "image72.svg"},
    "services": {"png": "image85.png", "svg": "image86.svg"},
    "devops": {"png": "image79.png", "svg": "image80.svg"},
    "datacenter": {"png": "image77.png", "svg": "image78.svg"},
    "networking": {"png": "image83.png", "svg": "image84.svg"},
    "spotlight": {"png": "image87.png", "svg": "image88.svg"},
    "wireless": {"png": "image101.png", "svg": "image102.svg"},
    "collaboration": {"png": "image89.png", "svg": "image90.svg"},
    "download": {"png": "image91.png", "svg": "image92.svg"},
    "systems": {"png": "image103.png", "svg": "image104.svg"},
    "managed_services": {"png": "image95.png", "svg": "image96.svg"},
    "multi_cloud": {"png": "image99.png", "svg": "image100.svg"},
    "location": {"png": "image93.png", "svg": "image94.svg"},
    "managers": {"png": "image97.png", "svg": "image98.svg"},
    "folder": {"png": "image113.png", "svg": "image114.svg"},
    "firewall": {"png": "image111.png", "svg": "image112.svg"},
    "employees": {"png": "image109.png", "svg": "image110.svg"},
    "email": {"png": "image107.png", "svg": "image108.svg"},
    "computer": {"png": "image105.png", "svg": "image106.svg"},
    "global": {"png": "image81.png", "svg": "image82.svg"},
    "data_analytics": {"png": "image75.png", "svg": "image76.svg"},
    # Slide 28 — Brand Icon Library page 2
    "vsoc": {"png": "image147.png", "svg": "image148.svg"},
    "modern_device_management": {"png": "image139.png", "svg": "image140.svg"},
    "dvi": {"png": "image129.png", "svg": "image130.svg"},
    "data": {"png": "image127.png", "svg": "image128.svg"},
    "data_security": {"png": "image125.png", "svg": "image126.svg"},
    "software_automation": {"png": "image141.png", "svg": "image142.svg"},
    "computer_link": {"png": "image121.png", "svg": "image122.svg"},
    "computer_chip": {"png": "image119.png", "svg": "image120.svg"},
    "incident_response": {"png": "image133.png", "svg": "image134.svg"},
    "cloud_management": {"png": "image117.png", "svg": "image118.svg"},
    "upload": {"png": "image145.png", "svg": "image146.svg"},
    "handshake": {"png": "image131.png", "svg": "image132.svg"},
    "consulting_services": {"png": "image123.png", "svg": "image124.svg"},
    "messaging": {"png": "image137.png", "svg": "image138.svg"},
    "iot_ot": {"png": "image135.png", "svg": "image136.svg"},
    "threat_detection": {"png": "image143.png", "svg": "image144.svg"},
    "security_analytics": {"png": "image157.png", "svg": "image158.svg"},
    "threat_hunting": {"png": "image159.png", "svg": "image160.svg"},
    "threat_intelligence": {"png": "image161.png", "svg": "image162.svg"},
    "laptop": {"png": "image155.png", "svg": "image156.svg"},
    "apps_infrastructure": {"png": "image149.png", "svg": "image150.svg"},
    "integrate_public_clouds": {"png": "image153.png", "svg": "image154.svg"},
    "devops_cicd": {"png": "image151.png", "svg": "image152.svg"},
    "transform_networking": {"png": "image163.png", "svg": "image164.svg"},
    "modern_workplace": {"png": "image175.png", "svg": "image176.svg"},
    "empower_digital_workspaces": {"png": "image169.png", "svg": "image170.svg"},
    "lifecycle_management": {"png": "image171.png", "svg": "image172.svg"},
    "software_defined_infrastructure": {"png": "image179.png", "svg": "image180.svg"},
    "digital_workspace": {"png": "image167.png", "svg": "image168.svg"},
    "servers": {"png": "image177.png", "svg": "image178.svg"},
    "desktop_transformation": {"png": "image165.png", "svg": "image166.svg"},
    "managed_services_alt": {"png": "image173.png", "svg": "image174.svg"},
    # Slide 29 — Brand Icon Library page 3
    "all_flash_storage": {"png": "image183.png", "svg": "image184.svg"},
    "resource": {"png": "image193.png", "svg": "image194.svg"},
    "full_stack_solutions": {"png": "image185.png", "svg": "image186.svg"},
    "protect": {"png": "image191.png", "svg": "image192.svg"},
    "hyper_converged": {"png": "image187.png", "svg": "image188.svg"},
    "virtual_desktop": {"png": "image195.png", "svg": "image196.svg"},
    "access": {"png": "image181.png", "svg": "image182.svg"},
    "modernize_data_center": {"png": "image189.png", "svg": "image190.svg"},
    "process": {"png": "image203.png", "svg": "image204.svg"},
    "infrastructure": {"png": "image201.png", "svg": "image202.svg"},
    "data_alt": {"png": "image199.png", "svg": "image200.svg"},
    "app": {"png": "image197.png", "svg": "image198.svg"},
    "predictable_expense": {"png": "image213.png", "svg": "image214.svg"},
    "variable_opex": {"png": "image219.png", "svg": "image220.svg"},
    "utility_billing": {"png": "image217.png", "svg": "image218.svg"},
    "metering_agreement": {"png": "image211.png", "svg": "image212.svg"},
    "identity": {"png": "image209.png", "svg": "image210.svg"},
    "scorpion_secops": {"png": "image215.png", "svg": "image216.svg"},
    "digital_solution": {"png": "image207.png", "svg": "image208.svg"},
    "zero_trust_protect": {"png": "image223.png", "svg": "image224.svg"},
    "zero_trust_visibility": {"png": "image229.png", "svg": "image230.svg"},
    "zero_trust_architecture": {"png": "image221.png", "svg": "image222.svg"},
    "zero_trust_policy": {"png": "image227.png", "svg": "image228.svg"},
    "zero_trust_monitor": {"png": "image225.png", "svg": "image226.svg"},
    "mobile_device": {"png": "image205.png", "svg": "image206.svg"},
    "ppm_services": {"png": "image231.png", "svg": "image232.svg"},
    "secure": {"png": "image233.png", "svg": "image234.svg"},
    "unified": {"png": "image235.png", "svg": "image236.svg"},
}

# Preset icon positions (in inches from top-left)
ICON_POSITIONS = {
    "top_right": {"x": Inches(11.5), "y": Inches(0.2), "size": Inches(1.0)},
    "top_left": {"x": Inches(0.3), "y": Inches(0.2), "size": Inches(1.0)},
    "before_title": {"x": Inches(0.3), "y": Inches(1.55), "size": Inches(0.8)},
}


# ─── Speaker Notes ───

def inject_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide. Newlines create separate paragraphs."""
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    # Clear default content
    tf.clear()
    paragraphs = notes_text.split('\n')
    for i, para_text in enumerate(paragraphs):
        if i == 0:
            tf.paragraphs[0].text = para_text
        else:
            p = tf.add_paragraph()
            p.text = para_text


# ─── Chart Injection ───

def inject_chart(slide, chart_spec, left, top, width, height):
    """
    Add a chart to the slide at the specified position.

    chart_spec: {
        "type": "column|bar|line|pie|doughnut|stacked_column|stacked_bar|area",
        "categories": ["Cat1", "Cat2", ...],
        "series": [{"name": "Series1", "values": [1, 2, ...]}, ...],
        "title": "Optional chart title",
        "show_legend": true,
        "show_data_labels": false,
        "show_gridlines": true
    }
    """
    chart_type_str = chart_spec.get("type", "column").lower()
    xl_chart_type = CHART_TYPE_MAP.get(chart_type_str)
    if xl_chart_type is None:
        print(f"  WARNING: Unknown chart type '{chart_type_str}', defaulting to column")
        xl_chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED

    chart_data = CategoryChartData()
    chart_data.categories = chart_spec.get("categories", [])

    for series in chart_spec.get("series", []):
        chart_data.add_series(series.get("name", ""), series.get("values", []))

    chart_frame = slide.shapes.add_chart(
        xl_chart_type, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    _style_chart(chart, chart_spec, chart_type_str)
    return chart_frame


def _style_chart(chart, chart_spec, chart_type_str):
    """Apply <company_name> brand styling to a chart."""
    # Chart title
    chart_title = chart_spec.get("title")
    if chart_title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = chart_title
        for run in chart.chart_title.text_frame.paragraphs[0].runs:
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.name = "Helvetica Neue"
            run.font.color.rgb = RGBColor(0x00, 0x44, 0x7A)
    else:
        chart.has_title = False

    # Legend
    show_legend = chart_spec.get("show_legend", True)
    if show_legend and chart_type_str not in ("pie", "doughnut"):
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
        chart.legend.font.name = "Helvetica Neue"
    elif chart_type_str in ("pie", "doughnut"):
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
        chart.legend.font.name = "Helvetica Neue"
    else:
        chart.has_legend = False

    # Apply brand colors to series/points
    is_pie_type = chart_type_str in ("pie", "doughnut")

    if is_pie_type:
        # For pie/doughnut, color each data point
        plot = chart.plots[0]
        series = plot.series[0]
        for i in range(len(chart_spec.get("categories", []))):
            point = series.points[i]
            color = <company_name>_CHART_COLORS[i % len(<company_name>_CHART_COLORS)]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = color
        # Data labels with category + percentage
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(10)
        data_labels.font.name = "Helvetica Neue"
        data_labels.show_category_name = True
        data_labels.show_percentage = True
        data_labels.show_value = False
    else:
        # For category charts, color each series
        plot = chart.plots[0]
        for i, series in enumerate(plot.series):
            color = <company_name>_CHART_COLORS[i % len(<company_name>_CHART_COLORS)]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
            # Line charts: set line color too
            if chart_type_str in ("line", "area"):
                series.format.line.color.rgb = color
                series.format.line.width = Pt(2.5)

        # Data labels
        if chart_spec.get("show_data_labels", False):
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.size = Pt(9)
            data_labels.font.name = "Helvetica Neue"
            data_labels.number_format = '0'
            data_labels.show_value = True

        # Axes styling
        try:
            category_axis = chart.category_axis
            category_axis.tick_labels.font.size = Pt(10)
            category_axis.tick_labels.font.name = "Helvetica Neue"
            category_axis.has_major_gridlines = False
        except Exception:
            pass

        try:
            value_axis = chart.value_axis
            value_axis.tick_labels.font.size = Pt(10)
            value_axis.tick_labels.font.name = "Helvetica Neue"
            value_axis.has_major_gridlines = chart_spec.get("show_gridlines", True)
            if value_axis.has_major_gridlines:
                value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
        except Exception:
            pass


# ─── Table Injection ───

def inject_table(slide, table_spec, left, top, width, height):
    """
    Add a styled table to the slide.

    table_spec: {
        "headers": ["Col1", "Col2", ...],
        "rows": [["val1", "val2", ...], ...],
        "header_color": "#0080BA",
        "stripe_color": "#F0F6FA",
        "font_size": 12,
        "col_widths": [3.0, 3.0, ...]  (inches, optional)
    }
    """
    headers = table_spec.get("headers", [])
    rows = table_spec.get("rows", [])
    if not headers:
        print("  WARNING: Table has no headers, skipping")
        return None

    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)
    font_size = Pt(table_spec.get("font_size", 12))
    header_color = _parse_hex_color(table_spec.get("header_color", "#0080BA"))
    stripe_color = _parse_hex_color(table_spec.get("stripe_color", "#F0F6FA"))

    # Calculate table height based on row count (not full body area)
    row_height = Inches(0.40)
    header_height = Inches(0.45)
    actual_height = header_height + row_height * len(rows)
    # Cap to available space but don't stretch
    table_height = min(actual_height, height)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_height)
    table = table_shape.table

    # Disable built-in banding (we'll style manually)
    table.first_row = False
    table.horz_banding = False

    # Column widths
    col_widths = table_spec.get("col_widths")
    if col_widths and len(col_widths) == n_cols:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    else:
        # Distribute evenly
        col_w = int(width / n_cols)
        for i in range(n_cols):
            table.columns[i].width = col_w

    # Header row
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        _style_table_cell(cell, font_size, bold=True, font_color=RGBColor(0xFF, 0xFF, 0xFF), fill_color=header_color)

    # Data rows
    for r_idx, row_data in enumerate(rows):
        for c_idx, cell_text in enumerate(row_data):
            if c_idx >= n_cols:
                break
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(cell_text)
            fill = stripe_color if r_idx % 2 == 0 else None
            _style_table_cell(cell, font_size, bold=False, font_color=RGBColor(0x33, 0x33, 0x33), fill_color=fill)

    return table_shape


def _style_table_cell(cell, font_size, bold=False, font_color=None, fill_color=None):
    """Apply styling to a single table cell."""
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = font_size
        paragraph.font.name = "Helvetica Neue"
        paragraph.font.bold = bold
        if font_color:
            paragraph.font.color.rgb = font_color
    cell.text_frame.word_wrap = True
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


# ─── KPI Card Injection ───

def inject_kpi_cards(slide, kpis, body_area):
    """
    Add 2-4 KPI metric cards to the slide.

    kpis: [{"value": "$4.2M", "label": "Revenue", "color": "#0080BA"}, ...]
    body_area: dict with left, top, width, height
    """
    n = len(kpis)
    if n < 1:
        return

    gap = Inches(0.3)
    total_gap = gap * (n - 1)
    card_width = int((body_area["width"] - total_gap) / n)
    card_height = Inches(3.6)
    card_top = body_area["top"] + int((body_area["height"] - card_height) / 2)

    for i, kpi in enumerate(kpis):
        card_left = body_area["left"] + i * (card_width + gap)
        color = _parse_hex_color(kpi.get("color", "#0080BA"))
        _add_single_kpi_card(slide, kpi, card_left, card_top, card_width, card_height, color)


def _add_single_kpi_card(slide, kpi, x, y, w, h, color):
    """Add a single KPI rounded-rectangle card."""
    from pptx.oxml.ns import qn
    from lxml import etree

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # No border

    # Set up text frame
    tf = shape.text_frame
    tf.word_wrap = True

    # Value paragraph (large)
    p_value = tf.paragraphs[0]
    p_value.alignment = PP_ALIGN.CENTER
    run_value = p_value.add_run()
    run_value.text = kpi.get("value", "")
    run_value.font.size = Pt(40)
    run_value.font.bold = True
    run_value.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run_value.font.name = "Helvetica Neue"

    # Label paragraph (smaller)
    p_label = tf.add_paragraph()
    p_label.alignment = PP_ALIGN.CENTER
    run_label = p_label.add_run()
    run_label.text = kpi.get("label", "")
    run_label.font.size = Pt(14)
    run_label.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run_label.font.name = "Helvetica Neue"

    # Vertical centering via XML
    txBody = tf._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 'ctr')


# ─── Timeline Injection ───

def inject_timeline(slide, milestones, body_area):
    """
    Add a horizontal timeline to the slide.

    milestones: [{"label": "Phase 1", "description": "Planning", "color": "#0080BA"}, ...]
    body_area: dict with left, top, width, height
    """
    n = len(milestones)
    if n < 1:
        return

    line_y = body_area["top"] + int(body_area["height"] * 0.50)
    line_left = body_area["left"] + Inches(0.8)
    line_right = body_area["left"] + body_area["width"] - Inches(0.8)
    line_width = line_right - line_left
    slide_right_edge = body_area["left"] + body_area["width"]

    # Draw horizontal line
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        line_left, line_y, line_right, line_y
    )
    connector.line.color.rgb = RGBColor(0x98, 0x99, 0x98)
    connector.line.width = Pt(3)

    # Calculate node positions
    node_radius = Inches(0.25)
    spacing = line_width / (n - 1) if n > 1 else 0

    for i, ms in enumerate(milestones):
        color = _parse_hex_color(ms.get("color", "#0080BA"))

        if n == 1:
            node_x = line_left + int(line_width / 2) - node_radius
        else:
            node_x = line_left + int(spacing * i) - node_radius

        node_y = line_y - node_radius

        # Circle node
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, node_x, node_y,
            node_radius * 2, node_radius * 2
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        # Label above node
        label_width = Inches(2.0)
        label_x = node_x + node_radius - int(label_width / 2)
        # Clamp to keep textbox within body area
        label_x = max(body_area["left"], min(label_x, slide_right_edge - label_width))
        label_y = node_y - Inches(0.9)
        label_box = slide.shapes.add_textbox(label_x, label_y, label_width, Inches(0.5))
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = ms.get("label", "")
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = color
        run.font.name = "Helvetica Neue"

        # Description below node
        desc_text = ms.get("description", "")
        if desc_text:
            desc_x = node_x + node_radius - int(label_width / 2)
            desc_x = max(body_area["left"], min(desc_x, slide_right_edge - label_width))
            desc_y = node_y + node_radius * 2 + Inches(0.15)
            desc_box = slide.shapes.add_textbox(desc_x, desc_y, label_width, Inches(0.7))
            tf2 = desc_box.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = desc_text
            run2.font.size = Pt(11)
            run2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
            run2.font.name = "Helvetica Neue"


# ─── Icon Injection (using template-embedded icons) ───

def _capture_template_icon_assets(prs):
    """
    Capture all icon image parts (PNG + SVG) from the template before slides
    are removed. Returns a dict mapping media filenames to (blob, content_type).
    """
    icon_assets = {}  # filename -> (blob, content_type)

    # Collect all media filenames we need
    needed_files = set()
    for icon_info in ICON_CATALOG.values():
        needed_files.add(icon_info["png"])
        needed_files.add(icon_info["svg"])

    # Scan all relationships across all slides and slide layouts
    parts_to_scan = []
    for slide in prs.slides:
        parts_to_scan.append(slide.part)
    for layout in prs.slide_layouts:
        parts_to_scan.append(layout.part)

    for part in parts_to_scan:
        for rel in part.rels.values():
            try:
                target = rel.target_part
                if hasattr(target, 'partname') and hasattr(target, 'blob'):
                    partname = str(target.partname)
                    filename = partname.split('/')[-1]
                    if filename in needed_files and filename not in icon_assets:
                        icon_assets[filename] = (target.blob, target.content_type)
            except Exception:
                continue

    return icon_assets


def inject_icons(slide, icons_spec, icon_assets):
    """
    Add icons to a slide from the template's embedded icon library.

    icons_spec: [
        {"name": "cloud", "position": "top_right"},
        {"name": "cybersecurity", "position": {"x": 10.0, "y": 0.5}, "size": 1.2}
    ]
    icon_assets: dict from _capture_template_icon_assets()
    """
    if not icons_spec or not icon_assets:
        return

    for icon_spec in icons_spec:
        icon_name = icon_spec.get("name", "").lower().replace(" ", "_").replace("-", "_")
        if icon_name not in ICON_CATALOG:
            print(f"  WARNING: Unknown icon '{icon_name}', skipping. Available: {', '.join(sorted(ICON_CATALOG.keys())[:10])}...")
            continue

        catalog_entry = ICON_CATALOG[icon_name]
        png_filename = catalog_entry["png"]

        if png_filename not in icon_assets:
            print(f"  WARNING: Icon image '{png_filename}' not found in template assets, skipping")
            continue

        blob, content_type = icon_assets[png_filename]

        # Determine position
        position = icon_spec.get("position", "top_right")
        if isinstance(position, str):
            pos = ICON_POSITIONS.get(position, ICON_POSITIONS["top_right"])
            x, y, size = pos["x"], pos["y"], pos["size"]
        elif isinstance(position, dict):
            x = Inches(position.get("x", 11.5))
            y = Inches(position.get("y", 0.2))
            size = Inches(icon_spec.get("size", 1.0))
        else:
            x, y, size = ICON_POSITIONS["top_right"]["x"], ICON_POSITIONS["top_right"]["y"], ICON_POSITIONS["top_right"]["size"]

        # Override size if specified at top level
        if "size" in icon_spec and not isinstance(position, dict):
            size = Inches(icon_spec["size"])

        # Add picture shape (python-pptx handles media embedding)
        pic = slide.shapes.add_picture(io.BytesIO(blob), x, y, size, size)


# ─── Main Injection Dispatch ───

def inject_slide(prs, slide_spec):
    """Add a single slide to the presentation based on the spec."""
    slide_type = slide_spec.get("type", "").upper()

    if slide_type not in LAYOUT_CATALOG:
        print(f"  WARNING: Unknown slide type '{slide_type}', skipping")
        return None

    catalog = LAYOUT_CATALOG[slide_type]
    layout_index = catalog["layout_index"]
    layout = prs.slide_layouts[layout_index]

    # Add the slide
    slide = prs.slides.add_slide(layout)

    phs = catalog["placeholders"]

    # Set title
    if "title" in phs and "title" in slide_spec:
        ph_idx = phs["title"]["idx"]
        try:
            ph = slide.placeholders[ph_idx]
            set_placeholder_text(ph, slide_spec["title"])
        except KeyError:
            print(f"  WARNING: Title placeholder idx={ph_idx} not found on {slide_type}")

    # Set subtitle
    if "subtitle" in phs and "subtitle" in slide_spec:
        ph_idx = phs["subtitle"]["idx"]
        try:
            ph = slide.placeholders[ph_idx]
            set_placeholder_text(ph, slide_spec["subtitle"])
        except KeyError:
            print(f"  WARNING: Subtitle placeholder idx={ph_idx} not found on {slide_type}")

    # Set body bullets
    if "body" in phs and "bullets" in slide_spec:
        ph_idx = phs["body"]["idx"]
        try:
            ph = slide.placeholders[ph_idx]
            set_bullets(ph, slide_spec["bullets"])
        except KeyError:
            print(f"  WARNING: Body placeholder idx={ph_idx} not found on {slide_type}")

    # Set body text (non-bullet, plain text)
    if "body" in phs and "body_text" in slide_spec:
        ph_idx = phs["body"]["idx"]
        try:
            ph = slide.placeholders[ph_idx]
            set_placeholder_text(ph, slide_spec["body_text"])
        except KeyError:
            print(f"  WARNING: Body placeholder idx={ph_idx} not found on {slide_type}")

    # Set two-column content
    if "left" in phs and "left_bullets" in slide_spec:
        ph_idx = phs["left"]["idx"]
        try:
            ph = slide.placeholders[ph_idx]
            set_bullets(ph, slide_spec["left_bullets"])
        except KeyError:
            print(f"  WARNING: Left placeholder idx={ph_idx} not found on {slide_type}")

    if "right" in phs and "right_bullets" in slide_spec:
        ph_idx = phs["right"]["idx"]
        try:
            ph = slide.placeholders[ph_idx]
            set_bullets(ph, slide_spec["right_bullets"])
        except KeyError:
            print(f"  WARNING: Right placeholder idx={ph_idx} not found on {slide_type}")

    # Set quote text (uses title placeholder)
    if slide_type == "QUOTE" and "quote" in slide_spec:
        ph_idx = phs["title"]["idx"]
        try:
            ph = slide.placeholders[ph_idx]
            set_placeholder_text(ph, slide_spec["quote"])
        except KeyError:
            print(f"  WARNING: Quote placeholder idx={ph_idx} not found on {slide_type}")

    # ─── New slide type dispatch ───

    # CHART: full-width chart replacing body
    if slide_type == "CHART" and "chart" in slide_spec:
        _remove_body_placeholder(slide)
        inject_chart(slide, slide_spec["chart"],
                     BODY_AREA["left"], BODY_AREA["top"],
                     BODY_AREA["width"], BODY_AREA["height"])

    # TABLE: full-width table replacing body
    elif slide_type == "TABLE" and "table" in slide_spec:
        _remove_body_placeholder(slide)
        inject_table(slide, slide_spec["table"],
                     BODY_AREA["left"], BODY_AREA["top"],
                     BODY_AREA["width"], BODY_AREA["height"])

    # CONTENT_CHART: bullets left, chart right
    elif slide_type == "CONTENT_CHART":
        if "bullets" in slide_spec:
            # Resize body placeholder to left half
            try:
                ph = slide.placeholders[1]
                ph.left = SPLIT_LEFT["left"]
                ph.top = SPLIT_LEFT["top"]
                ph.width = SPLIT_LEFT["width"]
                ph.height = SPLIT_LEFT["height"]
                set_bullets(ph, slide_spec["bullets"])
            except KeyError:
                print(f"  WARNING: Body placeholder not found for CONTENT_CHART bullets")
        else:
            _remove_body_placeholder(slide)
        if "chart" in slide_spec:
            inject_chart(slide, slide_spec["chart"],
                         SPLIT_RIGHT["left"], SPLIT_RIGHT["top"],
                         SPLIT_RIGHT["width"], SPLIT_RIGHT["height"])

    # CONTENT_TABLE: bullets left, table right
    elif slide_type == "CONTENT_TABLE":
        if "bullets" in slide_spec:
            try:
                ph = slide.placeholders[1]
                ph.left = SPLIT_LEFT["left"]
                ph.top = SPLIT_LEFT["top"]
                ph.width = SPLIT_LEFT["width"]
                ph.height = SPLIT_LEFT["height"]
                set_bullets(ph, slide_spec["bullets"])
            except KeyError:
                print(f"  WARNING: Body placeholder not found for CONTENT_TABLE bullets")
        else:
            _remove_body_placeholder(slide)
        if "table" in slide_spec:
            inject_table(slide, slide_spec["table"],
                         SPLIT_RIGHT["left"], SPLIT_RIGHT["top"],
                         SPLIT_RIGHT["width"], SPLIT_RIGHT["height"])

    # KPI: metric cards
    elif slide_type == "KPI" and "kpis" in slide_spec:
        _remove_body_placeholder(slide)
        inject_kpi_cards(slide, slide_spec["kpis"], BODY_AREA)

    # TIMELINE: horizontal timeline
    elif slide_type == "TIMELINE" and "milestones" in slide_spec:
        _remove_body_placeholder(slide)
        inject_timeline(slide, slide_spec["milestones"], BODY_AREA)

    # CONTENT with optional chart overlay (replaces bullets with full chart)
    elif slide_type == "CONTENT" and "chart" in slide_spec and "bullets" not in slide_spec:
        _remove_body_placeholder(slide)
        inject_chart(slide, slide_spec["chart"],
                     BODY_AREA["left"], BODY_AREA["top"],
                     BODY_AREA["width"], BODY_AREA["height"])

    # TWO_COLUMN with optional chart replacing a column
    elif slide_type == "TWO_COLUMN":
        if "left_chart" in slide_spec:
            _remove_placeholder_by_idx(slide, 1)
            inject_chart(slide, slide_spec["left_chart"],
                         SPLIT_LEFT["left"], SPLIT_LEFT["top"],
                         SPLIT_LEFT["width"], SPLIT_LEFT["height"])
        if "right_chart" in slide_spec:
            _remove_placeholder_by_idx(slide, 2)
            inject_chart(slide, slide_spec["right_chart"],
                         SPLIT_RIGHT["left"], SPLIT_RIGHT["top"],
                         SPLIT_RIGHT["width"], SPLIT_RIGHT["height"])

    return slide


def _capture_closing_slide_assets(prs):
    """
    Capture shapes and image relationships from the template's CLOSING slide
    (slide index 21) before slides are removed. The CLOSING layout has no
    placeholders — all visible content (Thank You text, contact info, orange
    line, <company_name> logo) lives as shapes on the original slide.
    """
    from pptx.oxml.ns import qn

    closing_layout_name = LAYOUT_CATALOG["CLOSING"]["layout_name"]

    # Find the template slide that uses the CLOSING layout
    closing_slide = None
    for slide in prs.slides:
        if slide.slide_layout.name == closing_layout_name:
            closing_slide = slide
            break

    if closing_slide is None:
        print("  WARNING: Could not find CLOSING template slide for asset capture")
        return None

    # Capture shape XML elements and any image relationships
    spTree = closing_slide._element.find(qn('p:cSld')).find(qn('p:spTree'))
    shape_elements = []
    image_parts = {}  # rId -> (blob, content_type)

    for elem in spTree:
        tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
        # Skip the group shape properties (nvGrpSpPr, grpSpPr) — these are
        # structural elements of the spTree, not content shapes
        if tag in ('nvGrpSpPr', 'grpSpPr'):
            continue

        shape_xml = copy.deepcopy(elem)
        shape_elements.append(shape_xml)

        # If this is a picture, capture the image relationship
        if tag == 'pic':
            blipFill = elem.find(qn('p:blipFill'))
            if blipFill is not None:
                blip = blipFill.find(qn('a:blip'))
                if blip is not None:
                    rEmbed = blip.get(qn('r:embed'))
                    if rEmbed:
                        try:
                            rel = closing_slide.part.rels[rEmbed]
                            image_parts[rEmbed] = (
                                rel.target_part.blob,
                                rel.target_part.content_type,
                            )
                        except KeyError:
                            pass

    return {"shapes": shape_elements, "images": image_parts}


def _inject_closing_shapes(slide, closing_assets):
    """Copy the captured CLOSING slide shapes onto a new slide."""
    from pptx.oxml.ns import qn
    from pptx.parts.image import ImagePart
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    if closing_assets is None:
        return

    spTree = slide._element.find(qn('p:cSld')).find(qn('p:spTree'))

    for shape_xml in closing_assets["shapes"]:
        new_shape = copy.deepcopy(shape_xml)
        tag = new_shape.tag.split('}')[-1] if '}' in new_shape.tag else new_shape.tag

        # For pictures, we need to re-create the image relationship on the new slide
        if tag == 'pic':
            blipFill = new_shape.find(qn('p:blipFill'))
            if blipFill is not None:
                blip = blipFill.find(qn('a:blip'))
                if blip is not None:
                    old_rId = blip.get(qn('r:embed'))
                    if old_rId and old_rId in closing_assets["images"]:
                        blob, content_type = closing_assets["images"][old_rId]
                        # Create an ImagePart and relate it to the new slide
                        from pptx.opc.packuri import PackURI
                        partname = PackURI('/ppt/media/closing_logo.png')
                        image_part = ImagePart.load(
                            partname, content_type,
                            slide.part.package, blob,
                        )
                        new_rId = slide.part.relate_to(image_part, RT.IMAGE)
                        blip.set(qn('r:embed'), new_rId)

        spTree.append(new_shape)


def generate_presentation(plan, template_path, output_path):
    """Generate a full presentation from a content plan."""
    template_path = Path(template_path)
    output_path = Path(output_path)

    if not template_path.exists():
        print(f"ERROR: Template not found: {template_path}")
        sys.exit(1)

    # Validate
    warnings = validate_plan(plan)
    if warnings:
        print(f"\n⚠ VALIDATION WARNINGS ({len(warnings)}):")
        for w in warnings:
            print(f"  - {w}")
        print()

    # Load template
    prs = Presentation(str(template_path))

    # Capture CLOSING slide assets before removing slides
    # (The CLOSING layout has no placeholders — shapes live on the slide itself)
    closing_assets = _capture_closing_slide_assets(prs)

    # Capture icon assets from template media library before removing slides
    icon_assets = _capture_template_icon_assets(prs)

    # Remove all existing slides from the template
    # We want a clean presentation using only the layouts
    from pptx.oxml.ns import qn
    prs_elem = prs.part._element  # the <p:presentation> XML element
    sldIdLst = prs_elem.find(qn('p:sldIdLst'))
    if sldIdLst is not None:
        sld_ids_to_remove = list(sldIdLst)
        for sldId in sld_ids_to_remove:
            rId = sldId.get(qn('r:id'))
            if rId:
                try:
                    prs.part.drop_rel(rId)
                except Exception:
                    pass
            sldIdLst.remove(sldId)

    print(f"Generating presentation with {len(plan['slides'])} slides...")

    # Inject slides
    for i, slide_spec in enumerate(plan["slides"], 1):
        slide_type = slide_spec.get("type", "?").upper()
        title = slide_spec.get("title", slide_spec.get("quote", ""))[:40]
        print(f"  Slide {i}: {slide_type} — \"{title}\"")
        slide = inject_slide(prs, slide_spec)

        if slide is not None:
            # For CLOSING slides, copy the captured shapes
            if slide_type == "CLOSING":
                _inject_closing_shapes(slide, closing_assets)

            # Inject speaker notes (available on any slide type)
            if "speaker_notes" in slide_spec:
                inject_speaker_notes(slide, slide_spec["speaker_notes"])

            # Inject icons (available on any slide type)
            if "icons" in slide_spec:
                inject_icons(slide, slide_spec["icons"], icon_assets)

            # Inject slide transition
            if _TRANSITIONS_AVAILABLE:
                transition_spec = slide_spec.get("transition")
                # Fall back to the recommended default for this slide type
                if transition_spec is None:
                    transition_spec = RECOMMENDED_TRANSITION.get(slide_type)
                if transition_spec:
                    inject_transition(slide, transition_spec, slide_type)

            # Inject element animations
            if _ANIMATIONS_AVAILABLE:
                animations_spec = slide_spec.get("animations")
                # None signals "use smart defaults"; empty list [] means "no animations"
                if "animations" not in slide_spec:
                    inject_animations(slide, None, slide_type)
                elif animations_spec:
                    inject_animations(slide, animations_spec, slide_type)

    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"\nSaved: {output_path}")
    return output_path


# ─── CLI ───

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Generate a <company_name> presentation from a content plan.')
    parser.add_argument('plan', help='Path to content plan JSON file')
    parser.add_argument('--template', '-t',
                        default=str(Path(__file__).parent.parent / 'assets' / '<company_name>_2021_template.pptx'),
                        help='Path to <company_name> template .pptx')
    parser.add_argument('--output', '-o',
                        default=str(Path(__file__).parent.parent / 'output' / 'generated.pptx'),
                        help='Output path for generated .pptx')

    args = parser.parse_args()

    with open(args.plan) as f:
        plan = json.load(f)

    generate_presentation(plan, args.template, args.output)
